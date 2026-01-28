
import streamlit as st
import sys
import os
import io
import datetime

# === 1. ページ設定 (必ず一番最初) ===
st.set_page_config(page_title="Biz PDF Converter Pro", layout="centered")

# ===========================
# テーマ設定関数
# ===========================
def apply_theme(theme):
    base_css = """
        img { border: 1px solid #ddd; border-radius: 5px; }
        .streamlit-expanderHeader { font-weight: bold; font-size: 1.2em; background-color: #f0f2f6; border-radius: 5px; }
        .stMarkdown h3 { border-bottom: 2px solid #ddd; padding-bottom: 5px; margin-top: 20px; }
    """
    
    css = ""
    if theme == "ビジネス (通常)":
        color_primary = "#4CAF50"
        text_color = "#2E7D32"
        css = f"""
            .stButton>button {{ background-color: white; color: {text_color}; border: 2px solid {color_primary}; border-radius: 5px; font-weight: bold; width: 100%; }}
            .stButton>button:hover {{ background-color: {color_primary}; color: white; }}
            h1 {{ color: {text_color}; }}
            .stProgress .st-bo {{ background-color: {color_primary}; }}
        """
    elif theme == "箱推し (全員)":
        css = f"""
            .stApp {{ background: linear-gradient(135deg, #fff0f0 25%, #fffff0 25%, #fffff0 50%, #fff0f5 50%, #fff0f5 75%, #f8f0ff 75%); }}
            .stButton>button {{ 
                background: linear-gradient(90deg, #E60033, #FFF100, #FF69B4, #800080); 
                color: white; border: none; border-radius: 20px; font-weight: bold; width: 100%; text-shadow: 1px 1px 2px black;
            }}
            .stButton>button:hover {{ opacity: 0.9; }}
            h1 {{ 
                background: linear-gradient(90deg, #E60033, #F2C000, #FF69B4, #800080);
                -webkit-background-clip: text;
                -webkit-text-fill-color: transparent;
            }}
            .stProgress .st-bo {{ background: linear-gradient(90deg, #E60033, #FFF100, #FF69B4, #800080); }}
        """
    else:
        color_primary = "#4CAF50"
        css = f"""
            .stButton>button {{ background-color: white; color: #2E7D32; border: 2px solid {color_primary}; border-radius: 5px; font-weight: bold; width: 100%; }}
            .stButton>button:hover {{ background-color: {color_primary}; color: white; }}
            h1 {{ color: #2E7D32; }}
            .stProgress .st-bo {{ background-color: {color_primary}; }}
        """

    st.markdown(f"<style>{base_css} {css}</style>", unsafe_allow_html=True)

# === 2. ライブラリ読み込み ===
try:
    import fitz  # PyMuPDF
    import cv2
    import numpy as np
    import pytesseract
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches, Emu, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError as e:
    st.error(f"ライブラリ不足: {e}")
    st.stop()

# ===========================
# サイドバー設定
# ===========================
st.sidebar.header("🎨 テーマ & テンプレート")
selected_theme = st.sidebar.selectbox("テーマカラー", ["ビジネス (通常)", "箱推し (全員)"])
apply_theme(selected_theme)

template_file = st.sidebar.file_uploader("会社のPPTXテンプレート (任意)", type="pptx", help="会社のロゴ入りスライドなどを適用したい場合にアップロードしてください。")

st.sidebar.divider()
st.sidebar.header("📄 出力設定")
slide_sizing = st.sidebar.radio("スライドサイズ", ["PDFに合わせる (推奨)", "16:9 (ワイド)", "4:3 (標準)"])
quality_mode = st.sidebar.select_slider("画質設定", options=["軽量", "標準", "高画質"], value="標準")

if quality_mode == "軽量": zoom_factor=1.0; jpeg_quality=70
elif quality_mode == "標準": zoom_factor=1.5; jpeg_quality=80
else: zoom_factor=2.0; jpeg_quality=95

st.sidebar.divider()
st.sidebar.header("🛡️ 加工・編集")
pdf_password = st.sidebar.text_input("PDFパスワード (必要な場合)", type="password")

footer_text = st.sidebar.text_input("フッター文字", placeholder="例：© 2024 My Company")
watermark_text = st.sidebar.text_input("透かし文字", value="")
use_patch = st.sidebar.checkbox("修正用パッチを配置", value=True)

st.sidebar.subheader("✂️ ロゴ消し")
use_erase = st.sidebar.checkbox("ロゴ/不要領域の白塗り", value=True)
erase_w = st.sidebar.slider("右端カット (px)", 0, 800, 350)
erase_h = st.sidebar.slider("下端カット (px)", 0, 500, 180)

st.sidebar.divider()
ocr_enabled = st.sidebar.checkbox("テキスト抽出 (ノートへ)", value=True)

# ===========================
# メイン画面
# ===========================
st.title("🏢 Biz PDF Converter Pro")
st.caption("パスワード解除・テンプレート適用に対応した、企業導入モデルです。")

# ★★★ 親切丁寧な使い方ガイド（フロー図＋全機能網羅） ★★★
with st.expander("🔰 初めての方はこちら（目的・使い方・NotebookLM活用）", expanded=False):
    tab0, tab1, tab2, tab3 = st.tabs(["📖 はじめに・目的", "① 基本の使い方", "② Pro機能・テクニック", "③ 文字の直し方"])
    
    with tab0:
        st.markdown("""
        ### 何のためのアプリなの？
        このアプリは、**「PDF資料を、自由に編集できるパワーポイントに変換する」**ための道具です。
        
        最近流行りのAI（**NotebookLM**など）で作った資料は、内容が良くても「PDF」でしか保存できないことが多く、**「ページの順番を変えたい」「会社指定の表紙をつけたい」**といった時に困ってしまいます。
        このアプリを使えば、そんなPDFを一瞬でパワーポイント化し、自由に加工できる素材に変えることができます。
        
        ### 🚀 活用フロー（NotebookLMの例）
        
        **STEP 1：AIで資料を作る**
        [NotebookLM] などに指示を出して、調査レポートやまとめ資料を作らせ、PDFで保存します。
        
        　　⬇︎ （PDFファイルをこのアプリに入れる）
        
        **STEP 2：パワポに変換する（★このアプリの出番）**
        このアプリを使って、PDFをパワーポイント形式（.pptx）に変換します。
        レイアウトが崩れないよう「画像」として貼り付けられます。
        
        　　⬇︎ （ダウンロードしたファイルを開く）
        
        **STEP 3：仕上げる**
        パワーポイントで開き、以下の作業をして完成です！
        * **合体:** 会社の定型フォーマット（表紙など）と合体させる。
        * **修正:** 直したい文字の上に「白い箱」を置いて修正する。
        * **印刷:** ページ番号などを整えて印刷・配布。
        """)

    with tab1:
        st.markdown("""
        ### まずはここから！ 3ステップで変換
        
        **1. ファイルを入れる**
        下の「Drag and drop file here」の枠に、PDFファイルを置いてください。
        * 📁 フォルダからマウスで引っ張ってきてもOKです。
        * 複数のファイルを一度に入れても大丈夫です（自動で1つのパワポにまとまります）。
        
        **2. プレビューを確認**
        ファイルを入れると、画面の真ん中に「画像のプレビュー」が出ます。
        * **赤枠の部分**が、白く塗りつぶされて消える場所です。
        * 左メニューの「✂️ ロゴ消し」のつまみを動かして、消したい場所（ページ番号など）を調整してください。
        
        **3. 変換スタート**
        準備ができたら、ピンクや緑色の**「変換スタート」ボタン**を押してください。
        しばらく待つと、「完了しました！」と出て、ダウンロードボタンが表示されます。
        """)
        
    with tab2:
        st.markdown("""
        ### Pro版の高度な機能
        左側のメニュー（サイドバー）で、いろいろな調整ができます。
        
        **🔐 パスワード付きPDFを変換したい**
        * 給与明細や契約書など、パスワードがかかっているファイルの場合、左メニューの**「PDFパスワード」**に入力してください。自動で解除して変換します。
        
        **🎨 会社のテンプレートを使いたい**
        * 左メニューの一番上にある**「会社のPPTXテンプレート」**に、いつも使っている表紙やロゴ入りのパワポファイルをアップロードしてください。そのデザインの上にPDFを貼り付けます。
        
        **🏢 会社名やコピーライトを入れたい**
        * **「フッター文字」**に会社名を入れると、全ページの左下に自動で入ります。
        * **「透かし文字」**に「社外秘」などと入れると、ページの中央に薄くスタンプされます。
        """)
        
    with tab3:
        st.markdown("""
        ### パワポにした後、文字を直したい！
        このアプリは、レイアウト崩れを防ぐために「ページ全体を画像」として貼り付けます。
        そのため、文字を直接カチカチして打ち直すことはできません。
        
        **その代わり、「修正用パッチ」機能を使います！**
        
        1. 変換後のパワポを開くと、右上に**「白い箱（テキストボックス）」**が置いてあります。
        2. 直したい文字の上に、その白い箱をマウスで移動させて、上から被せます。
        3. 箱の中に、正しい文字を入力してください。
        
        **💡 ヒント：元の文章データはどこ？**
        パワポの画面の下にある**「ノート」**という欄を見てください。
        AIが読み取った文章がそこに全部入っていますので、そこからコピーして使うと便利です。
        """)

# ===========================
# 関数定義
# ===========================

def load_pdf_doc(file_bytes, password=""):
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        if doc.is_encrypted:
            if not doc.authenticate(password):
                st.error("🔒 パスワードが間違っているか、入力されていません。サイドバーで入力してください。")
                return None
        return doc
    except Exception as e:
        return None

def preprocess_image_for_ocr(cv_img):
    gray = cv2.cvtColor(cv_img, cv2.COLOR_BGR2GRAY)
    _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    return binary

def adjust_image(cv_img, brightness=0, contrast=0):
    if brightness == 0 and contrast == 0: return cv_img
    alpha = (contrast + 100.0) / 100.0 
    beta = brightness
    return cv2.convertScaleAbs(cv_img, alpha=alpha, beta=beta)

def add_watermark(cv_img, text="CONFIDENTIAL"):
    h, w = cv_img.shape[:2]
    overlay = cv_img.copy()
    font = cv2.FONT_HERSHEY_SIMPLEX
    scale = w / 1000.0 * 2.0
    thickness = int(scale * 2)
    (text_w, text_h), _ = cv2.getTextSize(text, font, scale, thickness)
    x = int((w - text_w) / 2)
    y = int((h + text_h) / 2)
    cv2.putText(overlay, text, (x, y), font, scale, (150, 150, 150), thickness, cv2.LINE_AA)
    cv2.addWeighted(overlay, 0.2, cv_img, 0.8, 0, cv_img)
    return cv_img

# ===========================
# メイン処理
# ===========================
uploaded_files = st.file_uploader("PDFファイルをアップロード (複数可)", type="pdf", accept_multiple_files=True)

if uploaded_files:
    docs = []
    total_pages_all = 0
    
    for up_file in uploaded_files:
        file_bytes = up_file.read()
        doc = load_pdf_doc(file_bytes, pdf_password)
        if doc:
            docs.append((up_file.name, doc))
            total_pages_all += len(doc)
    
    if docs:
        first_doc = docs[0][1]
        
        # --- プレビュー ---
        st.divider()
        st.subheader("👁️ 仕上がりプレビュー")
        col1, col2 = st.columns([1, 2])
        with col1:
            preview_page_idx = st.number_input("確認ページ", min_value=1, max_value=len(first_doc), value=1) - 1
            if pdf_password and first_doc.is_encrypted:
                st.success("🔓 パスワード解除成功")
        with col2:
            page = first_doc[preview_page_idx]
            pix = page.get_pixmap(matrix=fitz.Matrix(1.0, 1.0))
            img_prev = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            cv_prev = cv2.cvtColor(np.array(img_prev), cv2.COLOR_RGB2BGR)
            h, w = cv_prev.shape[:2]
            
            if watermark_text: cv_prev = add_watermark(cv_prev, watermark_text)
            if use_erase:
                preview_display = cv_prev.copy()
                cv2.rectangle(preview_display, (w - erase_w, h - erase_h), (w, h), (0, 0, 255), 4)
                overlay = preview_display.copy()
                cv2.rectangle(overlay, (w - erase_w, h - erase_h), (w, h), (255, 200, 200), -1)
                cv2.addWeighted(overlay, 0.5, preview_display, 0.5, 0, preview_display)
            else:
                preview_display = cv_prev
            st.image(cv2.cvtColor(preview_display, cv2.COLOR_BGR2RGB), caption="プレビュー", use_container_width=True)

        # --- 変換実行 ---
        st.divider()
        btn_label = "Z伝説 変換スタート！" if selected_theme == "箱推し (全員)" else "変換スタート"
        
        if st.button(btn_label, type="primary"):
            p_bar = st.progress(0)
            status_area = st.empty()
            
            if template_file:
                try:
                    prs = Presentation(template_file)
                    layout = prs.slide_layouts[-1]
                except:
                    prs = Presentation()
                    layout = prs.slide_layouts[6]
            else:
                prs = Presentation()
                layout = prs.slide_layouts[6]
            
            page1 = first_doc[0]
            pdf_w, pdf_h = page1.rect.width, page1.rect.height
            
            if not template_file:
                if slide_sizing == "PDFに合わせる (推奨)":
                    prs.slide_width = Emu(pdf_w * 12700); prs.slide_height = Emu(pdf_h * 12700)
                elif slide_sizing == "16:9 (ワイド)":
                    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
                else:
                    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)

            current_cnt = 0
            for filename, doc in docs:
                status_area.text(f"処理中: {filename} ...")
                for i, page in enumerate(doc):
                    pix = page.get_pixmap(matrix=fitz.Matrix(zoom_factor, zoom_factor))
                    img_pil = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    cv_img = cv2.cvtColor(np.array(img_pil), cv2.COLOR_RGB2BGR)
                    h_orig, w_orig = cv_img.shape[:2]
                    
                    if watermark_text: cv_img = add_watermark(cv_img, watermark_text)
                    if use_erase:
                        cv2.rectangle(cv_img, (w_orig - int(erase_w * zoom_factor), h_orig - int(erase_h * zoom_factor)), 
                                      (w_orig, h_orig), (255, 255, 255), -1)
                    
                    slide = prs.slides.add_slide(layout)
                    img_bytes = cv2.imencode(".jpg", cv_img, [int(cv2.IMWRITE_JPEG_QUALITY), jpeg_quality])[1].tobytes()
                    image_stream = io.BytesIO(img_bytes)
                    slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                    
                    if use_patch:
                        patch_box = slide.shapes.add_textbox(prs.slide_width - Inches(2.5), Inches(0.2), Inches(2), Inches(0.5))
                        patch_box.fill.solid()
                        patch_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        tf_patch = patch_box.text_frame
                        tf_patch.text = "修正用ボックス"
                        tf_patch.paragraphs[0].font.size = Pt(10)
                        tf_patch.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

                    page_number_val = current_cnt + 1
                    txBox = slide.shapes.add_textbox(prs.slide_width - Inches(1.5), prs.slide_height - Inches(0.5), Inches(1), Inches(0.3))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = str(page_number_val)
                    p.font.size = Pt(12)
                    p.font.color.rgb = RGBColor(100, 100, 100)
                    p.alignment = PP_ALIGN.RIGHT
                    
                    if footer_text:
                        txBox2 = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.5), Inches(5), Inches(0.3))
                        tf2 = txBox2.text_frame
                        p2 = tf2.paragraphs[0]
                        p2.text = footer_text
                        p2.font.size = Pt(10)
                        p2.font.color.rgb = RGBColor(150, 150, 150)

                    if ocr_enabled:
                        try:
                            ocr_img = preprocess_image_for_ocr(cv_img)
                            header = f"[{filename} - P.{i+1}]\n"
                            text = pytesseract.image_to_string(ocr_img, lang='jpn+eng')
                            slide.notes_slide.notes_text_frame.text = header + text
                        except:
                            slide.notes_slide.notes_text_frame.text = ""
                    
                    current_cnt += 1
                    p_bar.progress(current_cnt / total_pages_all)
            
            success_msg = "完了しました！ ゼーーーット！" if selected_theme == "箱推し (全員)" else "変換が完了しました！"
            status_area.success(success_msg)
            
            out_ppt = io.BytesIO()
            prs.save(out_ppt)
            out_ppt.seek(0)
            
            today_str = datetime.datetime.now().strftime('%Y%m%d')
            base_name = os.path.splitext(docs[0][0])[0]
            dl_name = f"まとめ資料_変換済_{today_str}.pptx" if len(docs) > 1 else f"{base_name}_変換済_{today_str}.pptx"
            
            st.download_button("📥 パワポをダウンロード", out_ppt, dl_name)
